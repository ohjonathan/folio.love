"""Shared test fixtures for folio pipeline tests."""

import pytest
from pathlib import Path


@pytest.fixture
def tmp_output(tmp_path):
    """Provide a temporary output directory."""
    out = tmp_path / "output"
    out.mkdir()
    return out


@pytest.fixture
def sample_pptx(tmp_path):
    """Create a minimal valid PPTX file."""
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Slide"
    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def sample_pptx_multi(tmp_path):
    """Create a PPTX with 3 slides for boundary testing."""
    from pptx import Presentation
    prs = Presentation()
    for i in range(1, 4):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"Slide {i} Title"
    path = tmp_path / "multi.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def empty_file(tmp_path):
    """Create an empty file."""
    path = tmp_path / "empty.pptx"
    path.touch()
    return path


@pytest.fixture
def non_zip_pptx(tmp_path):
    """Create a .pptx that is not a ZIP file."""
    path = tmp_path / "fake.pptx"
    path.write_text("this is not a zip file")
    return path
